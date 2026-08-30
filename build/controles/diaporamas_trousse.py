#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Les diaporamas de la trousse : rien ne déborde, aucune image ne recouvre.

    python3 build/controles/diaporamas_trousse.py

N'écrit rien. Sort en **code 1** au premier écart, comme les autres contrôles
du dépôt, pour s'enchaîner dans un `&&`.

**Pourquoi il existe.** Le thème refuse un tableau trop haut, et c'est le seul
garde-fou de la chaîne. Tout le reste passe : une consigne de deux lignes dans
une boîte d'une ligne, une image posée sous un texte, un titre de carte qui
double et vient toucher son corps. Rien n'échoue, le fichier s'écrit, et le
défaut se voit pour la première fois **projeté sur un mur, devant une
direction**.

Trois défauts réels ont été trouvés par ce contrôle le 30 août 2026 : dix-neuf
consignes d'annexe à 0,675 pouce dans des boîtes de 0,62 (A1, A2), la consigne
du dialogue de F1, et le titre « Ce qui change pour l'enseignant » de P2.

**Deux pièges de mesure, tous deux payés en l'écrivant.**

  · Un titre se compose à 1,15 d'interligne, le corps à 1,35. Mesurer les
    titres comme du corps fait crier au débordement sur quatre diapositives
    parfaitement saines — et un contrôle qui a tort coûte plus cher que pas de
    contrôle, puisqu'on cherche une panne qui n'existe pas.
  · Une estimation maison du retour à la ligne (caractères ÷ largeur) rend zéro
    partout, y compris sur les défauts connus. La seule mesure qui vaut est
    `theme.h_of()`, celle que le thème emploie lui-même.
"""

import glob
import os
import sys

ICI = os.path.dirname(os.path.abspath(__file__))
RACINE = os.path.dirname(os.path.dirname(ICI))
PPT = os.path.join(RACINE, 'build', 'powerpoints')
sys.path.insert(0, PPT)

from pptx import Presentation          # noqa: E402
from pptx.util import Emu              # noqa: E402
from theme import FS, FOOT_Y, h_of     # noqa: E402

DIAPOS = os.path.join(RACINE, 'assets', 'presentations', 'diaporamas')
MARGE = 0.05        # tolérance de mesure, en pouces
SLIDE_W = 13.333


def _boites(sl):
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue
        runs = [r for p in sh.text_frame.paragraphs for r in p.runs if r.font.size]
        if not runs:
            continue
        yield (sh, txt, max(r.font.size.pt for r in runs),
               any(r.font.bold for r in runs))


def controler(chemin):
    """Rend la liste des écarts d'un diaporama."""
    ecarts = []
    nom = os.path.basename(chemin)
    for i, sl in enumerate(Presentation(chemin).slides, 1):
        for sh, txt, pts, gras in _boites(sl):
            w, h = Emu(sh.width).inches, Emu(sh.height).inches
            lh = 1.15 if (gras and pts >= FS['h3']) else 1.35
            besoin = h_of(txt, pts, gras, w, lh)
            if besoin > h + MARGE:
                ecarts.append('%s dia %2d : « %s… » demande %.2f po, la boîte en '
                              'fait %.2f' % (nom, i, txt.replace('\n', ' ')[:44],
                                             besoin, h))

        images = [s for s in sl.shapes if s.shape_type == 13]
        if not images:
            continue
        p = images[0]
        x, y = Emu(p.left).inches, Emu(p.top).inches
        w, h = Emu(p.width).inches, Emu(p.height).inches
        bas = max([Emu(s.top).inches + Emu(s.height).inches
                   for s, _, _, _ in _boites(sl) if Emu(s.top).inches < 4.0] or [0])
        if y < bas - 0.02:
            ecarts.append('%s dia %2d : l\'image monte sous le texte' % (nom, i))
        if y + h > FOOT_Y - 0.10:
            ecarts.append('%s dia %2d : l\'image mord la ligne de pied' % (nom, i))
        if x < 0.5 or x + w > SLIDE_W - 0.5:
            ecarts.append('%s dia %2d : l\'image sort de la gouttière' % (nom, i))
    return ecarts


def main():
    fichiers = sorted(glob.glob(os.path.join(DIAPOS, '*.pptx')))
    if not fichiers:
        print('!! aucun diaporama dans %s' % DIAPOS)
        return 1
    tous = []
    for f in fichiers:
        ec = controler(f)
        tous += ec
        print('  %-42s %s' % (os.path.basename(f),
                              'ok' if not ec else '%d écart(s)' % len(ec)))
    if tous:
        print()
        for e in tous:
            print('  !! %s' % e)
        print('\n%d écart(s) sur %d diaporama(s).' % (len(tous), len(fichiers)))
        return 1
    print('\n%d diaporama(s), aucun débordement.' % len(fichiers))
    return 0


if __name__ == '__main__':
    sys.exit(main())
