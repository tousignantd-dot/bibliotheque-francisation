"""
Socle visuel des présentations d'enseignement · Bibliothèque de francisation
============================================================================

Traduit `assets/design-system/` (jetons CSS) en gabarits PowerPoint.
Aucune valeur inventée : chaque couleur, chaque corps, chaque espacement
descend d'un jeton du système de design. Voir la table de correspondance
plus bas.

SUBSTITUTION SIGNALÉE (exigée par le README du système de design)
-----------------------------------------------------------------
Le système impose **Nunito**. PowerPoint n'embarque pas les polices sur
macOS et Nunito n'est pas installée sur le poste : un fichier livré avec
Nunito s'afficherait avec une police de remplacement arbitraire, et la
mise en page bougerait. On lui substitue **Verdana**, présente sur tout
poste Windows et macOS, choisie parce que c'est la police système la plus
proche des critères qui ont fait retenir Nunito : formes très ouvertes,
`a` et `g` à un seul étage, chasse large, excellente lisibilité projetée.
C'est le seul écart au système de design dans ce dossier.

ÉCHELLE TYPOGRAPHIQUE PROJETÉE
-------------------------------
Les corps du système visent un écran tenu à 40 cm. Une diapo se lit à
6 m du fond de la classe. On applique donc un facteur constant d'environ
1,6 à toute l'échelle, sans en changer les rapports :

    web 62 → 40 pt (titre de diapo)      web 19 → 20 pt (chapeau)
    web 30 → 28 pt (titre de bloc)       web 18 → 19 pt (énoncé)
    web 24 → 24 pt (mot vedette)         web 17 → 18 pt (consigne)
    web 13 → 13 pt (sur-titre capitales, seul corps inchangé)

Le plancher absolu est 14 pt : rien de plus petit n'est projeté, ce qui
transpose la règle « jamais sous 17 px » du système.
"""

import re
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import ImageFont

# ═══════════════════════════════════════════════════════════════════
#  1 · JETONS — copie fidèle de assets/design-system/tokens/colors.css
# ═══════════════════════════════════════════════════════════════════

C = {
    # Neutres chauds
    'paper_50':  'FBFBFA', 'paper_100': 'F7F7F5', 'paper_200': 'F0F0EE',
    'white':     'FFFFFF',
    'line_100':  'EAEAE8', 'line_200':  'D8E8DF', 'line_300':  'D6D6D2',
    'ink_900':   '17181A', 'ink_700':   '3A3D40',
    'ink_500':   '4B4F52', 'ink_400':   '6E7175',
    # Accent unique : le vert
    'green_800': '07734A', 'green_600': '0A8F5B',
    'green_100': 'E6F5EE', 'green_050': 'EDF6F1',
    # Accents de section — repérage seulement
    'acier_600': '1D6B8F', 'acier_100': 'E7F0F6',
    'indigo_600':'3B49A0', 'indigo_100':'E8EAFA',
    'ambre_700': 'B45309', 'ambre_100': 'FBEEDC',
    'teal_700':  '0D7A6F', 'teal_100':  'DCF2EF',
    'framboise_600': 'A5335F', 'framboise_100': 'FCE9F0',
    # Audio — le seul aplat rouge
    'audio':     'DC2626', 'audio_hover': 'B91C1C',
    # Rétroaction
    'ok_bg':   'E6F5EE', 'ok_line':   '0A8F5B', 'ok_ink':   '07734A',
    'no_bg':   'FFF6F5', 'no_line':   'DC2626', 'no_ink':   'B91C1C',
    'warn_bg': 'FEF6E7', 'warn_line': 'D9880B', 'warn_ink': '8A5206',
    # Sélection = plaque encre
    'sel_bg':  '17181A', 'sel_ink': 'FFFFFF',
}

# Les cinq couleurs de section, chacune avec sa teinte pâle. La graphie-phonie
# était violette ; le violet #6B4FBB est passé à la marque francis, qui le veut
# exclusif, et l'indigo l'a remplacé au repérage. Le vocabulaire était forêt ;
# le 20 août 2026 le vert est sorti du repérage — dans les modules comme ici,
# pour que la diapositive et l'écran s'accordent — et la framboise a pris sa
# place. Le seul vert qui reste est l'accent unique, `green_600`.
SECTIONS = {
    'acier':  ('acier_600',  'acier_100'),   # compréhension orale
    'indigo': ('indigo_600', 'indigo_100'),  # graphie-phonie
    'ambre':  ('ambre_700',  'ambre_100'),   # écriture
    'teal':   ('teal_700',   'teal_100'),    # écoute et réponds
    'framboise': ('framboise_600', 'framboise_100'),  # vocabulaire / bilan
}

FONT = 'Verdana'

# Corps projetés (voir en-tête)
FS = {
    'hero': 40, 'h2': 28, 'h3': 24, 'lead': 20,
    'body': 19, 'body_sm': 18, 'ui': 15, 'label': 13, 'meta': 12,
}
FS_MIN = 14   # plancher de projection

# Graisses du système : 400 / 600 / 700 / 800 / 900.
# Verdana n'a que régulier et gras — 700 et plus deviennent gras.

# Échelle d'espacement du système, en pouces (base 4 px ≈ 0,042")
SP = {1: 0.05, 2: 0.09, 3: 0.14, 4: 0.18, 5: 0.22, 6: 0.27, 7: 0.36, 8: 0.53}

RADIUS_CARD = 0.16   # 18 px transposé
RADIUS_CTRL = 0.10   # 10 px transposé

# Géométrie de la diapo : 16:9 large, format de projection.
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.62              # gouttière de page (32 px transposés)
CONTENT_W = SLIDE_W - 2 * MARGIN
HEAD_TOP = 0.46            # haut du sur-titre
BODY_TOP = 1.78            # haut de la zone de contenu des diapos courantes
FOOT_Y = 6.95              # ligne de pied

_FONT_CACHE = {}


# Mesure : on charge la police à 8× le corps demandé, puis on divise.
# À taille réelle, l'ajustement de grille (« hinting ») de PIL fausse la
# largeur de 2 à 4 % — assez pour qu'un titre calculé sur une ligne en
# occupe deux à l'écran. À 8×, l'erreur devient négligeable.
_MESURE = 8

# Marge de sécurité sur toute mesure. PowerPoint ne rend pas les glyphes
# exactement comme PIL ; 2 % absorbent l'écart sans se voir.
_SLACK = 1.02


def _pil(bold=False, size=20):
    key = (bold, size)
    if key not in _FONT_CACHE:
        path = ('/System/Library/Fonts/Supplemental/Verdana Bold.ttf' if bold
                else '/System/Library/Fonts/Supplemental/Verdana.ttf')
        _FONT_CACHE[key] = ImageFont.truetype(path, size)
    return _FONT_CACHE[key]


# Les balises d'emphase des séances, traduites en runs plutôt qu'imprimées.
#
# `theme.py` ne les a jamais interprétées : `<b>` partait tel quel dans le XML
# de la diapositive et s'affichait en clair au projecteur. Vérifié le 22 août
# 2026 sur les fichiers livrés — vingt-trois modules, plus de cent
# cinquante diapositives. Et le gras y est **pédagogique** : « le son de "un"
# ne s'écrit que <b>un</b> ou <b>um</b> ». Le supprimer aurait fait perdre
# exactement ce que la diapositive enseigne.
#
# **Liste blanche, jamais une expression générale.** Les séances contiennent
# du texte légitime entre chevrons — « <une réduction de loyer> », « <Le
# chauffe-eau du dessus> » — qu'un nettoyage naïf des balises avalerait. Seuls
# `<b>`, `<u>`, `<i>` et leurs fermetures sont reconnus ; tout le reste est du
# texte et le reste.
RE_EMPHASE = re.compile(r'</?(b|u|i)>', re.I)


def fragments(txt, gras=False, couleur='ink_900'):
    """Rend une liste de (texte, gras, couleur, souligné, italique).

    Un texte sans balise rend un seul fragment : le cas courant ne coûte rien.
    """
    if not txt or '<' not in txt:
        return [(txt, gras, couleur, False, False)]
    out, i = [], 0
    etat = {'b': gras, 'u': False, 'i': False}
    for m in RE_EMPHASE.finditer(txt):
        if m.start() > i:
            out.append((txt[i:m.start()], etat['b'], couleur, etat['u'], etat['i']))
        nom = m.group(1).lower()
        etat[nom] = not m.group(0).startswith('</')
        i = m.end()
    if i < len(txt):
        out.append((txt[i:], etat['b'], couleur, etat['u'], etat['i']))
    return [f for f in out if f[0]] or [('', gras, couleur, False, False)]


def sans_emphase(txt):
    """Le texte tel qu'il sera lu — pour mesurer, pas pour afficher."""
    return RE_EMPHASE.sub('', txt or '')


def text_w(txt, size, bold=False):
    """Largeur d'un texte, en pouces, marge de sécurité comprise."""
    px = _pil(bold, int(round(size * _MESURE))).getlength(txt)
    return px / _MESURE / 72.0 * _SLACK


def wrap(txt, size, bold, max_w):
    """Découpe un texte en lignes qui tiennent dans `max_w` pouces."""
    lines, cur = [], ''
    for word in txt.split():
        essai = word if not cur else cur + ' ' + word
        if text_w(essai, size, bold) <= max_w or not cur:
            cur = essai
        else:
            lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    return lines


def fit(txt, max_w, max_h, size, bold=False, lh=1.35, floor=FS_MIN):
    """Rend le plus grand corps ≤ `size` qui fait tenir `txt` dans la boîte.
    C'est le garde-fou contre les débordements : aucune diapo ne peut
    sortir d'ici avec du texte coupé."""
    s = size
    while s > floor:
        n = len(wrap(txt, s, bold, max_w))
        if n * s * lh / 72.0 <= max_h:
            return s
        s -= 1
    return floor


def h_of(txt, size, bold, max_w, lh=1.35):
    """Hauteur occupée par un texte, en pouces."""
    return len(wrap(txt, size, bold, max_w)) * size * lh / 72.0


# ═══════════════════════════════════════════════════════════════════
#  2 · PRIMITIVES DE DESSIN
# ═══════════════════════════════════════════════════════════════════

def _no_shadow(shape):
    """Le système bâtit la structure avec des filets, pas de la profondeur.
    python-pptx n'expose pas l'ombre : on la coupe dans le XML."""
    el = shape._element.spPr
    for tag in ('a:effectLst', 'a:effectRef'):
        for child in el.findall(qn(tag)):
            el.remove(child)
    el.append(el.makeelement(qn('a:effectLst'), {}))


def _spacing(run, em):
    """Interlettrage — le sur-titre en capitales est le seul usage
    (`--ls-label: .12em` dans le système)."""
    run.font._rPr.set('spc', str(int(em * run.font.size.pt * 100)))


class Slide:
    """Une diapo en cours de dessin. Toutes les coordonnées sont en pouces."""

    def __init__(self, deck, prs_slide):
        self.deck = deck
        self.s = prs_slide
        self.shapes = prs_slide.shapes

    # ── formes ──────────────────────────────────────────────────────
    def rect(self, x, y, w, h, fill=None, line=None, lw=1, radius=None):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sh = self.shapes.add_shape(shape_type, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
        if radius:
            # L'ajustement est une fraction de la moitié du petit côté.
            sh.adjustments[0] = min(0.5, radius / (min(w, h) / 2.0) / 2.0)
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = RGBColor.from_string(fill)
        else:
            sh.fill.background()
        if line:
            sh.line.color.rgb = RGBColor.from_string(line)
            sh.line.width = Pt(lw)
        else:
            sh.line.fill.background()
        sh.text_frame.text = ''
        _no_shadow(sh)
        return sh

    def text(self, txt, x, y, w, h, size=FS['body'], color='ink_900',
             bold=False, align='l', anchor='t', lh=1.35, italic=False,
             spc=None, autofit=True, lines=None):
        """Pose un bloc de texte. `lines` accepte une liste de tuples
        (texte, gras, couleur) pour mettre un fragment en évidence."""
        tb = self.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE,
                              'b': MSO_ANCHOR.BOTTOM}[anchor]
        if autofit and txt:
            # On mesure le texte LU, pas le texte écrit : « <b>um</b> » fait
            # deux caractères à l'écran, pas onze. Mesurer les balises faisait
            # rétrécir la police pour rien.
            size = fit(sans_emphase(txt), w, h, size, bold, lh)
        p = tf.paragraphs[0]
        p.alignment = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER,
                       'r': PP_ALIGN.RIGHT}[align]
        p.line_spacing = lh
        # Chaque segment est redécoupé sur ses balises d'emphase : un appel qui
        # passe déjà des `lines` en profite comme un appel simple.
        bruts = lines or [(txt, bold, color)]
        segments = []
        for seg in bruts:
            seg_txt, seg_bold, seg_col = seg[0], seg[1], seg[2]
            segments.extend(fragments(seg_txt, seg_bold, seg_col))
        for seg_txt, seg_bold, seg_col, seg_soul, seg_ital in segments:
            r = p.add_run()
            r.text = seg_txt
            r.font.size = Pt(size)
            r.font.bold = seg_bold
            r.font.italic = italic or seg_ital
            r.font.underline = seg_soul or None
            r.font.name = FONT
            r.font.color.rgb = RGBColor.from_string(C[seg_col])
            if spc:
                _spacing(r, spc)
        return tb

    def para_list(self, items, x, y, w, size=FS['body'], gap=0.16,
                  color='ink_700', lh=1.4):
        """Empile des paragraphes indépendants ; renvoie le bas atteint."""
        cy = y
        for it in items:
            hh = h_of(it, size, False, w, lh)
            self.text(it, x, cy, w, hh + 0.02, size=size, color=color,
                      lh=lh, autofit=False)
            cy += hh + gap
        return cy

    def image(self, path, x, y, w, h):
        """Insère une photo pédagogique CENTRÉE dans la boîte (x, y, w, h),
        sans jamais la déformer ni la laisser dépasser. Le système n'admet
        que des images pédagogiques, cadrées serré."""
        from PIL import Image as _Img
        with _Img.open(path) as im:
            ratio = im.width / im.height
        iw, ih = w, w / ratio
        if ih > h:
            ih, iw = h, h * ratio
        return self.shapes.add_picture(path, Inches(x + (w - iw) / 2),
                                       Inches(y + (h - ih) / 2),
                                       Inches(iw), Inches(ih))

    def notes(self, txt):
        self.s.notes_slide.notes_text_frame.text = txt


# ═══════════════════════════════════════════════════════════════════
#  3 · LE DECK ET SES GABARITS
# ═══════════════════════════════════════════════════════════════════

class Deck:
    """Une présentation. `section` choisit la couleur de repérage, qui
    n'apparaît qu'aux quatre endroits permis par le système : la pastille,
    le sur-titre, un filet gauche de 4 px, le point de parcours."""

    def __init__(self, code, section, titre, chapeau, duree):
        import modules
        m = modules.actif()
        self.MODULE, self.numero = m['titre'], m['numero']
        self.slug = modules.slug_actif()
        self.code, self.titre_deck = code, titre
        self.chapeau, self.duree = chapeau, duree
        self.sec, self.sec_soft = (C[a] for a in SECTIONS[section])
        self.sec_key, self.sec_soft_key = SECTIONS[section]
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    # ── socle commun ────────────────────────────────────────────────
    def _new(self, fond='paper_100', pied=True):
        sl = Slide(self, self.prs.slides.add_slide(self.blank))
        sl.rect(0, 0, SLIDE_W, SLIDE_H, fill=C[fond])
        self.n += 1
        if pied:
            sl.text(f'{self.code} · {self.MODULE}', MARGIN, FOOT_Y, 7.0, 0.24,
                    size=FS['meta'], color='ink_400', autofit=False)
            sl.text(str(self.n), SLIDE_W - MARGIN - 1.0, FOOT_Y, 1.0, 0.24,
                    size=FS['meta'], color='ink_400', align='r', autofit=False)
        return sl

    def _entete(self, sl, surtitre, titre):
        """Sur-titre coloré + titre encre. Le seul endroit où la couleur de
        section touche le haut de la diapo."""
        sl.text(surtitre.upper(), MARGIN, HEAD_TOP, CONTENT_W, 0.26,
                size=FS['label'], color=self.sec_key, bold=True,
                spc=0.12, autofit=False)
        size = fit(titre, CONTENT_W, 0.92, FS['h2'] + 6, True, 1.15)
        sl.text(titre, MARGIN, HEAD_TOP + 0.36, CONTENT_W, 0.96,
                size=size, color='ink_900', bold=True, lh=1.15, autofit=False)
        return sl

    def _carte(self, sl, x, y, w, h, teinte=None):
        return sl.rect(x, y, w, h, fill=teinte or C['white'],
                       line=C['line_100'], radius=RADIUS_CARD)

    # ── gabarit 1 · page de titre ───────────────────────────────────
    def titre(self, notes=''):
        """Bandeau clair pleine largeur — jamais noir, le système l'interdit."""
        sl = self._new(fond='green_050', pied=False)
        sl.text(f"SÉANCE {self.code}  ·  MODULE {self.numero}",
                MARGIN, 1.35, CONTENT_W, 0.3,
                size=FS['label'], color=self.sec_key, bold=True,
                spc=0.12, autofit=False)
        size = fit(self.titre_deck, CONTENT_W - 1.2, 2.0, FS['hero'], True, 1.05)
        sl.text(self.titre_deck, MARGIN, 1.85, CONTENT_W - 1.2, 2.05,
                size=size, color='ink_900', bold=True, lh=1.05, autofit=False)
        # La situation, dans une carte blanche — hiérarchie de page, niveau 1.
        hh = h_of(self.chapeau, FS['lead'], False, CONTENT_W - 1.9, 1.45)
        self._carte(sl, MARGIN, 4.15, CONTENT_W - 1.2, hh + 0.72)
        sl.text(self.chapeau, MARGIN + 0.36, 4.51, CONTENT_W - 1.92, hh + 0.05,
                size=FS['lead'], color='ink_500', lh=1.45, autofit=False)
        sl.text(f"{self.duree}  ·  {self.MODULE}", MARGIN, 6.6, CONTENT_W, 0.3,
                size=FS['ui'], color='ink_400', autofit=False)
        sl.notes(notes or f"Séance {self.code}. Durée prévue : {self.duree}.")
        return sl

    # ── gabarit 2 · objectifs ───────────────────────────────────────
    def objectifs(self, items, notes=''):
        """« À la fin, je serai capable de… » — la voix de l'apprenant,
        comme les titres d'exercice du module."""
        sl = self._new()
        self._entete(sl, 'Objectifs de la séance',
                     'À la fin, je serai capable de…')
        y, w = BODY_TOP + 0.15, CONTENT_W - 1.15
        for i, it in enumerate(items, 1):
            hh = max(0.62, h_of(it, FS['lead'], False, w - 0.5, 1.4) + 0.44)
            self._carte(sl, MARGIN, y, CONTENT_W, hh)
            # pastille numérotée teintée — usage permis de la couleur
            sl.rect(MARGIN + 0.26, y + (hh - 0.44) / 2, 0.44, 0.44,
                    fill=self.sec_soft, radius=0.14)
            sl.text(str(i), MARGIN + 0.26, y + (hh - 0.44) / 2 + 0.09, 0.44, 0.3,
                    size=FS['ui'], color=self.sec_key, bold=True, align='c',
                    autofit=False)
            sl.text(it, MARGIN + 0.92, y + 0.22, w - 0.5, hh - 0.4,
                    size=FS['lead'], color='ink_700', lh=1.4, autofit=False)
            y += hh + SP[3]
        sl.notes(notes or "Lire les objectifs à voix haute. Y revenir au billet de sortie.")
        return sl

    # ── gabarit 3 · déclencheur ─────────────────────────────────────
    def declencheur(self, surtitre, question, image=None, pistes=None,
                    notes=''):
        """Deux colonnes : la question à gauche, la photo pédagogique à droite.
        Aucune illustration décorative — c'est la règle du système."""
        sl = self._new()
        self._entete(sl, surtitre, question)
        if image:
            gauche_w = 7.1
            cx, cw = MARGIN + gauche_w + 0.4, CONTENT_W - gauche_w - 0.4
            ch = FOOT_Y - 0.3 - (BODY_TOP + 0.1)
            self._carte(sl, cx, BODY_TOP + 0.1, cw, ch)
            sl.image(image, cx + 0.24, BODY_TOP + 0.34, cw - 0.48, ch - 0.48)
        else:
            gauche_w = CONTENT_W
        if pistes:
            y = BODY_TOP + 0.25
            for p in pistes:
                hh = max(0.58, h_of(p, FS['body'], False, gauche_w - 0.95, 1.4) + 0.4)
                self._carte(sl, MARGIN, y, gauche_w, hh, teinte=C['white'])
                sl.rect(MARGIN, y + 0.12, 0.055, hh - 0.24, fill=self.sec)
                sl.text(p, MARGIN + 0.42, y + 0.2, gauche_w - 0.75, hh - 0.36,
                        size=FS['body'], color='ink_700', lh=1.4, autofit=False)
                y += hh + SP[2]
        sl.notes(notes or "Laisser le groupe répondre avant de donner quoi que ce soit.")
        return sl

    # ── gabarit 4 · la règle en une phrase ──────────────────────────
    def regle(self, surtitre, phrase, precision=None, notes=''):
        """Un énoncé, gros, seul. C'est la diapo qu'on photographie."""
        sl = self._new()
        sl.text(surtitre.upper(), MARGIN, HEAD_TOP, CONTENT_W, 0.26,
                size=FS['label'], color=self.sec_key, bold=True,
                spc=0.12, autofit=False)
        # Les deux cartes sont dimensionnées sur leur contenu et empilées :
        # aucune hauteur fixe, donc rien ne peut sortir de la diapo, quelle
        # que soit la longueur de la règle.
        top, bas = 1.34, FOOT_Y - 0.26
        pw, prw = CONTENT_W - 1.6, CONTENT_W - 0.9
        hp = h_of(precision, FS['lead'], False, prw, 1.45) + 0.62 if precision else 0
        size = 34
        while size > 20:
            h1 = h_of(phrase, size, True, pw, 1.22) + 1.0
            if top + h1 + (0.3 + hp if precision else 0) <= bas:
                break
            size -= 1
        h1 = h_of(phrase, size, True, pw, 1.22) + 1.0
        self._carte(sl, MARGIN, top, CONTENT_W, h1, teinte=C[self.sec_soft_key])
        sl.text(phrase, MARGIN + 0.8, top + 0.5, pw, h1 - 0.9,
                size=size, color='ink_900', bold=True, lh=1.22, autofit=False)
        if precision:
            py = top + h1 + 0.3
            self._carte(sl, MARGIN, py, CONTENT_W, hp)
            sl.text(precision, MARGIN + 0.44, py + 0.3, prw, hp - 0.5,
                    size=FS['lead'], color='ink_500', lh=1.45, autofit=False)
        sl.notes(notes or "Faire répéter la règle à voix haute par deux élèves.")
        return sl

    # ── gabarit 5 · tableau d'analyse ───────────────────────────────
    def tableau(self, surtitre, titre, entetes, lignes, note=None, notes='',
                cle=None, props=None):
        """Rangées séparées par un filet, jamais encadrées — hiérarchie
        de page, niveau 4. `cle` met une colonne en évidence (index)."""
        sl = self._new()
        self._entete(sl, surtitre, titre)
        ncol = len(entetes)
        y0 = BODY_TOP + 0.2
        bas = FOOT_Y - 0.26
        # Un tableau à deux colonnes est presque toujours un couple
        # « libellé · contenu » : la deuxième colonne reçoit le double de
        # largeur, sinon le contenu se replie sur trois lignes pour rien.
        props = props or ([0.34, 0.66] if ncol == 2 else [1 / ncol] * ncol)
        dispo_w = CONTENT_W - 0.8
        xs = [MARGIN + 0.4 + dispo_w * sum(props[:i]) for i in range(ncol)]
        ws = [dispo_w * p - 0.24 for p in props]
        notew = CONTENT_W - 0.9

        # Les rangées s'ajustent à leur contenu : une cellule qui passe à
        # deux lignes pousse toute la rangée, elle ne déborde jamais sur la
        # suivante. On cherche le plus grand corps qui fait tenir l'ensemble
        # — tableau ET note — dans la hauteur disponible.
        def mesurer(s, ns):
            hs = [max(0.5, max(h_of(str(c), s, cle is not None and i == cle,
                                    ws[i], 1.32)
                               for i, c in enumerate(lg)) + 0.26)
                  for lg in lignes]
            hn = (h_of(note, ns, False, notew, 1.45) + 0.56) if note else 0
            return hs, 0.5 + sum(hs) + 0.2, hn

        size, nsize = FS['body_sm'], FS['body_sm']
        while True:
            hs, haut, hn = mesurer(size, nsize)
            if y0 + haut + (0.26 + hn if note else 0) <= bas:
                break
            if note and nsize > FS_MIN:
                nsize -= 1
            elif size > FS_MIN:
                size -= 1
            else:
                # Au plancher et ça ne rentre toujours pas : c'est le contenu
                # qui est trop long, pas la mise en page. On refuse plutôt que
                # de livrer une diapo tronquée.
                raise ValueError(
                    f"{self.code} · tableau « {titre} » : trop de contenu pour "
                    "une diapositive projetée. Coupez-le en deux tableaux, "
                    "raccourcissez les cellules, ou retirez la note.")

        self._carte(sl, MARGIN, y0, CONTENT_W, haut)
        for i, e in enumerate(entetes):
            sl.text(e.upper(), xs[i], y0 + 0.18, ws[i], 0.26,
                    size=FS['label'], color='ink_400', bold=True, spc=0.12,
                    autofit=False)
        y = y0 + 0.5
        for li, ligne in enumerate(lignes):
            if li:
                sl.rect(MARGIN + 0.4, y, CONTENT_W - 0.8, 0.011, fill=C['line_100'])
            for i, cell in enumerate(ligne):
                est_cle = (cle is not None and i == cle)
                sl.text(str(cell), xs[i], y + 0.14, ws[i],
                        hs[li] - 0.22, size=size, bold=est_cle, lh=1.32,
                        color=(self.sec_key if est_cle else 'ink_700'),
                        autofit=False)
            y += hs[li]
        if note:
            ny = y0 + haut + 0.26
            self._carte(sl, MARGIN, ny, CONTENT_W, hn, teinte=C['paper_50'])
            sl.text(note, MARGIN + 0.44, ny + 0.28, notew, hn - 0.5,
                    size=nsize, color='ink_500', lh=1.45, autofit=False)
        sl.notes(notes or "Faire lire une colonne à la fois, jamais le tableau entier d'un coup.")
        return sl

    # ── gabarit 6 · cartes en grille ────────────────────────────────
    def cartes(self, surtitre, titre, items, cols=2, notes=''):
        """items = [(en-tête, corps)] — grille 2 ou 3 colonnes."""
        sl = self._new()
        self._entete(sl, surtitre, titre)
        rows = (len(items) + cols - 1) // cols
        gap = SP[4]
        cw = (CONTENT_W - gap * (cols - 1)) / cols
        dispo = FOOT_Y - 0.3 - (BODY_TOP + 0.2)
        ch = (dispo - gap * (rows - 1)) / rows
        # Un seul corps de titre et un seul corps de texte pour TOUTES les
        # cartes : le pire cas fixe la taille commune. Des cartes de corps
        # différents se voient immédiatement et cassent la grille.
        # Recherche conjointe : le titre de carte cède de la place au texte
        # avant que le texte ne descende sous le plancher. Pour un public en
        # apprentissage, c'est le corps de l'explication qui doit rester
        # grand, pas l'intitulé.
        tw = cw - 0.7
        meilleur = None
        for hs in range(FS['h3'], FS_MIN - 1, -1):
            htmax = max(h_of(h, hs, True, tw, 1.18) for h, _ in items)
            reste = ch - 0.3 - htmax - 0.36
            for ps in range(min(FS['body_sm'], hs - 4), FS_MIN - 1, -1):
                # Le titre de carte garde toujours 4 points d'avance sur le
                # texte : sans cet écart, la hiérarchie disparaît à l'écran.
                if max(h_of(p, ps, False, tw, 1.42) for _, p in items) <= reste:
                    if meilleur is None or (ps, hs) > (meilleur[1], meilleur[0]):
                        meilleur = (hs, ps, htmax, reste)
                    break
        if meilleur is None:   # cas désespéré : plancher partout
            hs = ps = FS_MIN
            htmax = max(h_of(h, hs, True, tw, 1.18) for h, _ in items)
            reste = ch - 0.3 - htmax - 0.36
        else:
            hs, ps, htmax, reste = meilleur
        for i, (h, p) in enumerate(items):
            cx = MARGIN + (i % cols) * (cw + gap)
            cy = BODY_TOP + 0.2 + (i // cols) * (ch + gap)
            self._carte(sl, cx, cy, cw, ch)
            sl.text(h, cx + 0.35, cy + 0.3, tw, htmax + 0.03,
                    size=hs, color=self.sec_key, bold=True, lh=1.18,
                    autofit=False)
            sl.text(p, cx + 0.35, cy + 0.36 + htmax, tw, reste,
                    size=ps, color='ink_700', lh=1.42, autofit=False)
        sl.notes(notes or "Une carte à la fois. Demander un exemple au groupe entre chaque.")
        return sl

    # ── gabarit 7 · dialogue ────────────────────────────────────────
    def dialogue(self, surtitre, titre, repliques, consigne=None, notes=''):
        """repliques = [(locuteur, texte, en_evidence)] — reprend
        `.dialogue__line` du système : le locuteur en gras, le texte à côté."""
        sl = self._new()
        self._entete(sl, surtitre, titre)
        y = BODY_TOP + 0.12
        if consigne:
            sl.text(consigne, MARGIN, y, CONTENT_W, 0.36, size=FS['body_sm'],
                    color='ink_400', autofit=False)
            y += 0.46
        dispo = FOOT_Y - 0.25 - y
        n = len(repliques)
        gap = SP[2]
        nomw = max([text_w(r[0], FS['body_sm'], True) for r in repliques]) + 0.3
        txtw = CONTENT_W - nomw - 1.0
        # corps commun à toutes les répliques, pour un bloc régulier
        size = FS['body']
        while size > FS_MIN:
            tot = sum(max(0.54, h_of(r[1], size, False, txtw, 1.4) + 0.36) + gap
                      for r in repliques) - gap
            if tot <= dispo:
                break
            size -= 1
        for qui, txt, evid in repliques:
            hh = max(0.54, h_of(txt, size, False, txtw, 1.4) + 0.36)
            if evid:
                self._carte(sl, MARGIN, y, CONTENT_W, hh,
                            teinte=C[self.sec_soft_key])
            sl.text(qui, MARGIN + 0.4, y + 0.18, nomw, hh - 0.3,
                    size=size, color=self.sec_key, bold=True, autofit=False)
            sl.text(txt, MARGIN + 0.4 + nomw, y + 0.18, txtw, hh - 0.3,
                    size=size, color='ink_700', lh=1.4, autofit=False)
            y += hh + gap
        sl.notes(notes or "Faire écouter l'audio d'abord, diapo masquée. Lire ensuite.")
        return sl

    # ── gabarit 8 · le piège ────────────────────────────────────────
    def piege(self, surtitre, faux, juste, explication, notes=''):
        """Deux colonnes : ce qu'on entend souvent, ce qu'il faut dire.
        Jamais la couleur seule — chaque colonne porte son mot (✕ / ✓)."""
        sl = self._new()
        self._entete(sl, surtitre, 'Le piège le plus fréquent')
        cw = (CONTENT_W - SP[5]) / 2
        y = BODY_TOP + 0.25
        # Verdana n'a ni ✓ ni ✕ : c'est le MOT qui double la couleur, ce que
        # le système autorise explicitement (« un glyphe OU un mot »).
        colonnes = [
            ('On entend souvent', faux, 'no_bg', 'no_line', 'no_ink'),
            ('Il faut dire plutôt', juste, 'ok_bg', 'ok_line', 'ok_ink'),
        ]
        for i, (mot, phrase, bg, ln, ink) in enumerate(colonnes):
            cx = MARGIN + i * (cw + SP[5])
            sl.rect(cx, y, cw, 2.55, fill=C[bg], line=C[ln], radius=RADIUS_CARD)
            sl.text(mot.upper(), cx + 0.4, y + 0.3, cw - 0.8, 0.3,
                    size=FS['label'], color=ink, bold=True, spc=0.12,
                    autofit=False)
            ps = fit(phrase, cw - 0.8, 1.5, 25, True, 1.3)
            sl.text(phrase, cx + 0.4, y + 0.78, cw - 0.8, 1.55,
                    size=ps, color='ink_900', bold=True, lh=1.3, autofit=False)
        ey = y + 2.85
        hh = h_of(explication, FS['lead'], False, CONTENT_W - 0.9, 1.45)
        self._carte(sl, MARGIN, ey, CONTENT_W, min(hh + 0.6, FOOT_Y - 0.3 - ey))
        sl.text(explication, MARGIN + 0.44, ey + 0.3, CONTENT_W - 0.9, hh + 0.05,
                size=FS['lead'], color='ink_500', lh=1.45, autofit=False)
        sl.notes(notes or "Demander qui a déjà fait cette erreur. Dédramatiser.")
        return sl

    # ── gabarit 9 · pratique guidée ─────────────────────────────────
    def pratique(self, surtitre, titre, consigne, items, corrige=None,
                 cols=1, notes=''):
        """Les exercices du module, projetés. `corrige` sert à la diapo
        suivante, générée automatiquement, avec la même mise en page."""
        for passe in (0, 1) if corrige else (0,):
            sl = self._new()
            self._entete(sl, surtitre if not passe else 'Correction', titre)
            sl.text(consigne if not passe else
                    "Correction — on lit la réponse, on ne commente pas l'erreur.",
                    MARGIN, BODY_TOP + 0.02, CONTENT_W, 0.36,
                    size=FS['body_sm'], color='ink_400', autofit=False)
            y0 = BODY_TOP + 0.56
            dispo = FOOT_Y - 0.28 - y0
            gap = SP[3]
            rows = (len(items) + cols - 1) // cols
            cw = (CONTENT_W - gap * (cols - 1)) / cols
            tw = cw - 1.05

            # Les deux passes doivent avoir EXACTEMENT la même géométrie —
            # sinon les rangées bougent au moment de la correction et le
            # groupe perd le fil. On dimensionne donc toujours sur le texte
            # corrigé, qui est le plus long des deux.
            longs = [it[0] + ('   » ' + it[1] if len(it) > 1 else '')
                     for it in items]

            def hauteurs(s):
                h = [max(0.56, h_of(t, s, False, tw, 1.4) + 0.34) for t in longs]
                if cols > 1:                       # grille : rangées égales
                    hmax = max(h)
                    h = [hmax] * len(h)
                    return h, hmax * rows + gap * (rows - 1)
                return h, sum(h) + gap * (len(h) - 1)

            size = FS['body']
            while size > FS_MIN and hauteurs(size)[1] > dispo:
                size -= 1
            hs, _ = hauteurs(size)

            for i, it in enumerate(items):
                enonce = it[0]
                rep = it[1] if len(it) > 1 else ''
                rh = hs[i]
                cx = MARGIN + (i % cols) * (cw + gap)
                cy = y0 + sum(hs[j] + gap for j in range(i)) if cols == 1 else \
                     y0 + (i // cols) * (rh + gap)
                self._carte(sl, cx, cy, cw, rh,
                            teinte=C['ok_bg'] if passe else C['white'])
                sl.rect(cx + 0.3, cy + (rh - 0.4) / 2, 0.4, 0.4,
                        fill=self.sec_soft, radius=0.13)
                sl.text(str(i + 1), cx + 0.3, cy + (rh - 0.4) / 2 + 0.07, 0.4, 0.28,
                        size=FS['ui'], color=self.sec_key, bold=True,
                        align='c', autofit=False)
                if passe and rep:
                    seg = [(enonce + '   ', False, 'ink_700'),
                           ('» ' + rep, True, 'ok_ink')]
                    txt = enonce + '   » ' + rep
                else:
                    seg, txt = None, enonce
                sl.text(txt, cx + 0.85, cy + 0.17, tw, rh - 0.3, size=size,
                        color='ink_700', lh=1.4, autofit=False, lines=seg)
            sl.notes(
                (notes or "Faire répondre le groupe avant d'afficher la correction. "
                          "Laisser du silence : les élèves lisent plus lentement qu'on ne "
                          "le croit.")
                if not passe else
                "Corriger en groupe. Donner la bonne réponse, sans revenir sur qui s'est trompé.")
        return self

    # ── gabarit 9b · l'exercice à l'écran ───────────────────────────
    def capture(self, exo, titre, consigne=None, notes=''):
        """L'écran même de l'exercice interactif, projeté avant qu'on l'ouvre.

        Elle vient APRÈS la diapositive « Pratique » qui reprend le même
        exercice en texte : le groupe a d'abord travaillé à l'oral, puis voit
        ce qu'il va manipuler. Elle ne la remplace pas — une image projetée ne
        se lit pas de la dernière rangée aussi bien qu'un texte composé pour
        elle, et l'enseignante a besoin du corrigé écrit.

        `exo` est l'identifiant de l'exercice dans le module (`pr1`), celui-là
        même que l'en-tête du deck cite. L'image est produite par
        `python3 captures.py <module>` ; on ne la fabrique pas ici, sinon
        chaque construction relancerait un navigateur.
        """
        import os
        chemin = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                              '_captures', self.slug, exo + '.png')
        if not os.path.exists(chemin):
            raise SystemExit(
                f'{self.code} : capture manquante pour « {exo} ».\n'
                f'   Attendue : {chemin}\n'
                f'   Produisez-la : python3 captures.py {self.slug}')
        sl = self._new()
        self._entete(sl, "Dans l'activité interactive", titre)
        y = BODY_TOP + 0.02
        sl.text(consigne or "L'écran que les élèves vont ouvrir. "
                            "Les réponses sont encore dans le banc, à droite.",
                MARGIN, y, CONTENT_W, 0.36, size=FS['body_sm'],
                color='ink_400', autofit=False)
        y += 0.5
        haut = FOOT_Y - 0.22 - y
        # Même calcul d'ajustement que `Slide.image`, mais mené ici parce que
        # le cadre doit épouser l'image : un filet posé sur la boîte pleine
        # flotterait autour d'une capture plus étroite qu'elle.
        from PIL import Image as _Img
        with _Img.open(chemin) as im:
            ratio = im.width / im.height
        iw, ih = CONTENT_W, CONTENT_W / ratio
        if ih > haut:
            ih, iw = haut, haut * ratio
        ix, iy = MARGIN + (CONTENT_W - iw) / 2, y + (haut - ih) / 2
        sl.rect(ix - 0.02, iy - 0.02, iw + 0.04, ih + 0.04,
                fill=C['white'], line=C['line_300'], radius=RADIUS_CTRL)
        sl.shapes.add_picture(chemin, Inches(ix), Inches(iy),
                              Inches(iw), Inches(ih))
        sl.notes(notes or
                 "Projeter avant d'ouvrir l'activité. Faire dire à deux ou "
                 "trois élèves où irait la première réponse, sans y toucher : "
                 "on vérifie qu'ils ont compris la consigne, pas le contenu. "
                 "Ouvrir ensuite le module sur les postes.")
        return sl

    # ── gabarit 10 · vocabulaire ────────────────────────────────────
    def vocabulaire(self, surtitre, titre, mots, notes=''):
        """mots = [(mot, définition)] — l'article fait partie du mot."""
        sl = self._new()
        self._entete(sl, surtitre, titre)
        y0 = BODY_TOP + 0.2
        haut = FOOT_Y - 0.3 - y0
        self._carte(sl, MARGIN, y0, CONTENT_W, haut)
        rowh = (haut - 0.36) / len(mots)
        # La colonne de gauche ne dépasse jamais 45 % de la largeur : sans
        # cette borne, une entrée longue (une phrase entière plutôt qu'un mot)
        # écrase la définition jusqu'à la faire disparaître de la diapo.
        motw = min(max(text_w(m, FS['body'], True) for m, _ in mots) + 0.45,
                   0.45 * (CONTENT_W - 0.8))
        y = y0 + 0.18
        for i, (mot, defi) in enumerate(mots):
            if i:
                sl.rect(MARGIN + 0.4, y, CONTENT_W - 0.8, 0.011,
                        fill=C['line_100'])
            ms = fit(mot, motw - 0.3, rowh - 0.2, FS['body'], True, 1.28)
            sl.text(mot, MARGIN + 0.4, y + 0.16, motw - 0.3, rowh - 0.24,
                    size=ms, color='ink_900', bold=True, lh=1.28, autofit=False)
            dw = CONTENT_W - 0.8 - motw
            ds = fit(defi, dw, rowh - 0.2, FS['body_sm'], False, 1.35)
            sl.text(defi, MARGIN + 0.4 + motw, y + 0.16, dw, rowh - 0.24,
                    size=ds, color='ink_500', lh=1.35, autofit=False)
            y += rowh
        sl.notes(notes or "Faire répéter chaque mot avec son article. L'article fait partie du mot.")
        return sl

    # ── gabarit 11 · billet de sortie ───────────────────────────────
    def billet(self, consigne, exemples=None, notes=''):
        """Le seul bloc foncé autorisé, et il est en fin de séance —
        le système en permet un par page, jamais l'en-tête."""
        sl = self._new()
        sl.rect(0, 0, SLIDE_W, SLIDE_H, fill=C['ink_900'])
        sl.text('BILLET DE SORTIE', MARGIN, 1.15, CONTENT_W, 0.3,
                size=FS['label'], color='white', bold=True, spc=0.12,
                autofit=False)
        size = fit(consigne, CONTENT_W - 1.4, 2.2, 34, True, 1.25)
        sl.text(consigne, MARGIN, 1.72, CONTENT_W - 1.4, 2.3,
                size=size, color='white', bold=True, lh=1.25, autofit=False)
        if exemples:
            y = 4.3
            for ex in exemples:
                hh = max(0.56, h_of(ex, FS['lead'], False, CONTENT_W - 1.1, 1.4) + 0.34)
                sl.rect(MARGIN, y, CONTENT_W, hh, fill='2A2C2F', radius=RADIUS_CTRL)
                sl.text(ex, MARGIN + 0.4, y + 0.17, CONTENT_W - 1.1, hh - 0.3,
                        size=FS['lead'], color='white', lh=1.4, autofit=False)
                y += hh + SP[2]
        sl.text(f'{self.code} · {self.MODULE}', MARGIN, FOOT_Y, 7.0, 0.24,
                size=FS['meta'], color='ink_400', autofit=False)
        sl.notes(notes or "Chacun écrit sa phrase et la remet en sortant. Deux minutes, pas plus.")
        return sl

    # ── enregistrement ──────────────────────────────────────────────
    def _verifier_glyphes(self):
        """Refuse d'enregistrer si un caractère n'existe pas dans Verdana.
        Sans ce contrôle, PowerPoint remplace le glyphe manquant par la
        police de secours du poste : ✓ et → deviennent un carré vide chez
        l'enseignant, et on ne le voit jamais depuis ici."""
        try:
            from fontTools.ttLib import TTFont
        except ImportError:
            return []
        dispo = set()
        for f in ('Verdana.ttf', 'Verdana Bold.ttf'):
            tt = TTFont('/System/Library/Fonts/Supplemental/' + f)
            for t in tt['cmap'].tables:
                dispo |= set(t.cmap.keys())
        manquants = {}
        for i, slide in enumerate(self.prs.slides, 1):
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        for ch in r.text:
                            if ord(ch) not in dispo and ch not in ' \n\t':
                                manquants.setdefault(ch, []).append(i)
        return [(ch, sorted(set(n))) for ch, n in manquants.items()]

    def save(self, dossier):
        import os
        mq = self._verifier_glyphes()
        if mq:
            détail = ' · '.join(f'{ch!r} (diapos {n})' for ch, n in mq)
            raise ValueError(
                f'{self.code} : caractères absents de Verdana — {détail}. '
                "Remplacez-les : ils s'afficheraient en carré vide.")
        nom = f'{self.code}-{_slug(self.titre_deck)}.pptx'
        chemin = os.path.join(dossier, nom)
        self.prs.save(chemin)
        return chemin, self.n


def _slug(s):
    import unicodedata, re
    # NFKD ne décompose pas les ligatures : « vœux » deviendrait « vux ».
    # On les écrit d'abord en toutes lettres.
    for lig, deux in (('œ', 'oe'), ('Œ', 'OE'), ('æ', 'ae'), ('Æ', 'AE')):
        s = s.replace(lig, deux)
    s = unicodedata.normalize('NFKD', s).encode('ascii', 'ignore').decode()
    s = re.sub(r"[^a-zA-Z0-9]+", '-', s).strip('-').lower()
    return re.sub(r'-+', '-', s)[:48]
