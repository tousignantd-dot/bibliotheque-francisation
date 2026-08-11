"""
Rendu d'épreuve · relit un .pptx et le redessine en PNG.

Le poste n'a pas LibreOffice. Ce script ouvre le fichier produit avec
python-pptx, lit la géométrie réelle de chaque forme et la repeint avec
Pillow en Verdana — la même police que celle déclarée dans le fichier.
Le rendu n'est donc pas une simulation de la mise en page : il lit ce qui
a été écrit dans le fichier.

    python apercu.py chemin/deck.pptx [dossier_sortie]
"""
import sys, os
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

PPP = 110  # points par pouce du rendu
VER = '/System/Library/Fonts/Supplemental/Verdana.ttf'
VERB = '/System/Library/Fonts/Supplemental/Verdana Bold.ttf'
_cache = {}


def f(bold, pt):
    k = (bold, round(pt))
    if k not in _cache:
        _cache[k] = ImageFont.truetype(VERB if bold else VER,
                                       max(1, int(round(pt * PPP / 72))))
    return _cache[k]


def po(v):
    return Emu(v).inches if v is not None else 0


def rgb(c):
    try:
        return '#' + str(c.rgb)
    except Exception:
        return None


def mes(txt, bold, pt):
    """Mesure haute résolution, identique à celle de theme.text_w, pour que
    l'épreuve casse les lignes exactement là où le générateur l'a prévu."""
    return _cache_m(bold, pt).getlength(txt) / 8.0 * PPP / 72.0


def _cache_m(bold, pt):
    k = ('m', bold, round(pt, 2))
    if k not in _cache:
        _cache[k] = ImageFont.truetype(VERB if bold else VER,
                                       int(round(pt * 8)))
    return _cache[k]


def wrap(dr, txt, font, maxpx, bold=False, pt=18):
    out, cur = [], ''
    for w in txt.split():
        t = w if not cur else cur + ' ' + w
        if mes(t, bold, pt) <= maxpx or not cur:
            cur = t
        else:
            out.append(cur); cur = w
    if cur:
        out.append(cur)
    return out


def render(path, outdir):
    prs = Presentation(path)
    W = int(po(prs.slide_width) * PPP)
    H = int(po(prs.slide_height) * PPP)
    os.makedirs(outdir, exist_ok=True)
    faits = []
    for i, slide in enumerate(prs.slides, 1):
        im = Image.new('RGB', (W, H), 'white')
        dr = ImageDraw.Draw(im)
        for sh in slide.shapes:
            x, y = po(sh.left) * PPP, po(sh.top) * PPP
            w, h = po(sh.width) * PPP, po(sh.height) * PPP
            if sh.shape_type is not None and sh.has_text_frame and sh.fill.type is not None:
                try:
                    fill = rgb(sh.fill.fore_color) if sh.fill.type == 1 else None
                except Exception:
                    fill = None
                try:
                    line = rgb(sh.line.color) if sh.line.fill.type == 1 else None
                except Exception:
                    line = None
                if fill or line:
                    dr.rounded_rectangle([x, y, x + w, y + h], radius=10,
                                         fill=fill, outline=line, width=1)
            if not sh.has_text_frame:
                if sh.shape_type == 13:  # image
                    dr.rectangle([x, y, x + w, y + h], outline='#B0B0B0')
                    dr.text((x + 8, y + 8), '[image]', fill='#909090',
                            font=f(False, 11))
                continue
            tf = sh.text_frame
            cy = y
            for p in tf.paragraphs:
                runs = [r for r in p.runs if r.text]
                if not runs:
                    continue
                pt = runs[0].font.size.pt if runs[0].font.size else 18
                bold = bool(runs[0].font.bold)
                col = rgb(runs[0].font.color) or '#17181A'
                lh = p.line_spacing or 1.35
                full = ''.join(r.text for r in runs)
                for ln in wrap(dr, full, f(bold, pt), max(4, w), bold, pt):
                    lw = mes(ln, bold, pt)
                    lx = x
                    if str(p.alignment) .startswith('CENTER'):
                        lx = x + (w - lw) / 2
                    elif str(p.alignment).startswith('RIGHT'):
                        lx = x + w - lw
                    # signale un débordement horizontal
                    if lw > w + 2:
                        dr.rectangle([lx, cy, lx + lw, cy + pt * PPP / 72],
                                     outline='#FF00FF', width=2)
                    dr.text((lx, cy), ln, fill=col, font=f(bold, pt))
                    cy += pt * lh * PPP / 72
                if cy > y + h + 3 and w > 2 and h > 2:
                    dr.rectangle([x, y, x + w, y + h], outline='#FF00FF', width=2)
        out = os.path.join(outdir, f'{i:02d}.png')
        im.save(out, quality=88)
        faits.append(out)
    return faits


if __name__ == '__main__':
    src = sys.argv[1]
    dest = sys.argv[2] if len(sys.argv) > 2 else 'apercu'
    for p in render(src, dest):
        print(p)
